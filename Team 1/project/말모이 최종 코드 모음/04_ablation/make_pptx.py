"""
말모이 팀 - 승정원일기 번역 전략 비교 발표 자료
python make_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

NAVY  = RGBColor(0x1a, 0x2a, 0x4a)
BLUE  = RGBColor(0x27, 0x5d, 0xa6)
GOLD  = RGBColor(0xd4, 0xa0, 0x17)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY  = RGBColor(0xf2, 0xf4, 0xf8)
GREEN = RGBColor(0x1a, 0x7a, 0x4a)
RED   = RGBColor(0xc0, 0x39, 0x2b)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]  # 완전 빈 슬라이드


def add_slide():
    return prs.slides.add_slide(blank)


def bg(slide, color):
    from pptx.util import Emu
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, text, font_size=24, bold=False,
        fg=WHITE, bg_color=NAVY, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = fg
    # 배경
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    sp = txBox._element
    spPr = sp.find(qn('p:spPr'))
    solidFill = etree.SubElement(spPr, qn('a:solidFill'))
    srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', f'{bg_color.rgb:06X}')
    return txBox


def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def text(slide, left, top, width, height, content, size=20, bold=False,
         color=NAVY, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = content
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def multiline(slide, left, top, width, height, lines, size=18, color=NAVY):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line_text, bold in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


# ─────────────────────────────────────────────
# 슬라이드 1: 표지
# ─────────────────────────────────────────────
s1 = add_slide()
bg(s1, NAVY)
rect(s1, 0, 5.8, 13.33, 1.7, BLUE)
rect(s1, 0, 0, 0.25, 7.5, GOLD)

text(s1, 0.5, 0.6, 12, 1.0, '말모이 팀', 28, bold=True, color=GOLD)
text(s1, 0.5, 1.5, 12, 1.8,
     '승정원일기 자동 번역 시스템\nNER 주입 기반 4-way 전략 비교',
     40, bold=True, color=WHITE)
text(s1, 0.5, 3.5, 10, 0.8,
     'Gemma-4-26B 기반  |  eval 300  |  BLEU · chrF · NER Recall',
     20, color=RGBColor(0xaa, 0xcc, 0xff))
text(s1, 0.5, 6.0, 12, 0.6, '2026', 18, color=WHITE)


# ─────────────────────────────────────────────
# 슬라이드 2: 프로젝트 개요
# ─────────────────────────────────────────────
s2 = add_slide()
bg(s2, GRAY)
rect(s2, 0, 0, 13.33, 1.3, NAVY)
rect(s2, 0, 0, 0.25, 7.5, GOLD)
text(s2, 0.5, 0.2, 12, 0.9, '프로젝트 개요', 32, bold=True, color=WHITE)

cards = [
    ('📜  데이터', '승정원일기\n조선시대 왕실 비서 일기\n한문 원문 → 현대 한국어 번역'),
    ('🎯  목표',   'LLM 기반 자동 번역 품질 향상\n인물명(NER) 정확도 개선\n4가지 프롬프트 전략 비교'),
    ('📊  평가셋', '전체 말뭉치 62,176개\n→ 계층적 샘플링 300개 추출\n(NER 포함 188개 / 미포함 112개)'),
    ('🤖  모델',   'Gemma-4-26B-A4B-IT\nGoogle Gemini API\ntemperature=0.0'),
]

for i, (title, body) in enumerate(cards):
    col = i % 2
    row = i // 2
    lft = 0.4 + col * 6.5
    tp  = 1.5 + row * 2.7
    rect(s2, lft, tp, 6.1, 2.4, WHITE)
    rect(s2, lft, tp, 6.1, 0.55, BLUE)
    text(s2, lft + 0.15, tp + 0.08, 5.8, 0.45, title, 18, bold=True, color=WHITE)
    text(s2, lft + 0.15, tp + 0.65, 5.8, 1.6, body, 16, color=NAVY)


# ─────────────────────────────────────────────
# 슬라이드 3: 평가셋 300개 구성
# ─────────────────────────────────────────────
s3 = add_slide()
bg(s3, GRAY)
rect(s3, 0, 0, 13.33, 1.3, NAVY)
rect(s3, 0, 0, 0.25, 7.5, GOLD)
text(s3, 0.5, 0.2, 12, 0.9, '평가셋 300개 구성', 32, bold=True, color=WHITE)

rect(s3, 0.4, 1.4, 5.8, 5.6, WHITE)
text(s3, 0.55, 1.55, 5.5, 0.5, '샘플링 방법', 20, bold=True, color=NAVY)
steps = [
    '① 전체 말뭉치 62,176개 확보',
    '② eval_set_1925: 1,925개 후보 추출',
    '   (NER 보유 기사 우선 포함)',
    '③ 1,925개 → 300개 계층 샘플링',
    '   · 원문 길이 분포 유지',
    '   · NER 엔티티 수 분포 유지',
    '④ reference 번역 포함된 항목만 선정',
]
for j, s in enumerate(steps):
    text(s3, 0.55, 2.1 + j * 0.6, 5.5, 0.55, s, 16,
         bold=s.startswith('①') or s.startswith('②') or s.startswith('③') or s.startswith('④'),
         color=NAVY)

rect(s3, 6.6, 1.4, 6.3, 2.5, WHITE)
text(s3, 6.75, 1.55, 6.0, 0.5, 'NER 구성', 20, bold=True, color=NAVY)
ner_data = [
    ('NER 엔티티 있음', '188개', '62.7%', BLUE),
    ('NER 엔티티 없음', '112개', '37.3%', RGBColor(0x88,0x99,0xbb)),
]
for j, (label, cnt, pct, c) in enumerate(ner_data):
    tp = 2.1 + j * 1.0
    rect(s3, 6.75, tp, pct_w := float(pct[:-1]) / 100 * 5.5, 0.45, c)
    text(s3, 6.75, tp, 3.5, 0.45, f' {label}', 15, color=WHITE)
    text(s3, 6.75 + 3.6, tp, 2.5, 0.45, f'{cnt}  ({pct})', 15, color=NAVY)

rect(s3, 6.6, 4.1, 6.3, 2.8, WHITE)
text(s3, 6.75, 4.25, 6.0, 0.5, '엔티티 수 분포', 20, bold=True, color=NAVY)
dist = [(0, 112), (1, 106), (2, 44), (3, 17), (4, 11), ('6+', 10)]
max_v = 112
for j, (k, v) in enumerate(dist):
    bar_w = v / max_v * 4.5
    rect(s3, 6.75, 4.85 + j * 0.32, bar_w, 0.25, BLUE)
    text(s3, 6.75, 4.85 + j * 0.32, 0.5, 0.28, str(k)+'개', 12, color=NAVY)
    text(s3, 6.75 + bar_w + 0.05, 4.85 + j * 0.32, 1.0, 0.28, str(v), 12, color=NAVY)


# ─────────────────────────────────────────────
# 슬라이드 4: 4가지 전략
# ─────────────────────────────────────────────
s4 = add_slide()
bg(s4, GRAY)
rect(s4, 0, 0, 13.33, 1.3, NAVY)
rect(s4, 0, 0, 0.25, 7.5, GOLD)
text(s4, 0.5, 0.2, 12, 0.9, '4가지 번역 전략', 32, bold=True, color=WHITE)

strategies = [
    ('① Baseline',
     '기본 번역 프롬프트',
     '별도 예시 없이 모델에게\n직접 한문 번역 요청',
     NAVY),
    ('② Few-shot',
     '번역 예시 5개 제공',
     '실록 문체 예시 5개를\n프롬프트에 포함하여 번역',
     BLUE),
    ('③ Few + NERfix',
     '예시 + 후처리 교정',
     '번역 후 인물 사전으로\nNER 오류를 후처리 교정',
     RGBColor(0x15, 0x6b, 0x8a)),
    ('④ NER Inject',
     '예시 + 인물명 직접 주입',
     '[등장 인물] 블록으로\n한자→한글 매핑을 프롬프트에 삽입',
     GREEN),
]

for i, (title, subtitle, desc, color) in enumerate(strategies):
    lft = 0.4 + i * 3.2
    rect(s4, lft, 1.4, 3.0, 5.5, WHITE)
    rect(s4, lft, 1.4, 3.0, 0.7, color)
    text(s4, lft + 0.1, 1.45, 2.8, 0.6, title, 18, bold=True, color=WHITE)
    text(s4, lft + 0.1, 2.2, 2.8, 0.5, subtitle, 16, bold=True, color=color)
    text(s4, lft + 0.1, 2.85, 2.8, 2.0, desc, 15, color=NAVY)

    # 프롬프트 예시 박스
    rect(s4, lft + 0.1, 4.95, 2.8, 1.8, GRAY)
    prompts = {
        '① Baseline':       '번역하세요:\n{원문}',
        '② Few-shot':       '[번역 예시]\n- 또 아뢰기를...\n{원문}',
        '③ Few + NERfix':   '[번역 예시]\n- ...\n→ 후처리 교정',
        '④ NER Inject':     '[등장 인물]\n· 金堉→김육\n위 문체로 번역:\n{원문}',
    }
    text(s4, lft + 0.15, 5.05, 2.7, 1.6, prompts[title], 11, color=RGBColor(0x33,0x33,0x33))


# ─────────────────────────────────────────────
# 슬라이드 5: 결과 비교
# ─────────────────────────────────────────────
s5 = add_slide()
bg(s5, GRAY)
rect(s5, 0, 0, 13.33, 1.3, NAVY)
rect(s5, 0, 0, 0.25, 7.5, GOLD)
text(s5, 0.5, 0.2, 12, 0.9, '실험 결과 비교  (n=291)', 32, bold=True, color=WHITE)

# 표
cols  = ['전략', 'BLEU(c)', 'chrF', 'BLEU(m)', 'NER Recall']
cw    = [3.2, 2.2, 2.2, 2.2, 2.2]
rows  = [
    ['Baseline',      '35.51', '30.72', '27.35', '0.803'],
    ['Few-shot',      '39.46', '34.27', '31.90', '0.792'],
    ['Few + NERfix',  '39.78', '34.69', '32.19', '0.938'],
    ['NER Inject ★', '40.61', '35.17', '33.09', '1.000'],
]

table_left = 0.5
table_top  = 1.5
row_h = 0.75

# 헤더
x = table_left
for ci, (col, cw_) in enumerate(zip(cols, cw)):
    rect(s5, x, table_top, cw_, row_h, NAVY)
    text(s5, x + 0.1, table_top + 0.15, cw_ - 0.1, row_h - 0.1,
         col, 17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += cw_

# 데이터 행
for ri, row in enumerate(rows):
    is_best = (ri == len(rows) - 1)
    row_color = RGBColor(0xe8, 0xf5, 0xec) if is_best else WHITE
    x = table_left
    for ci, (cell, cw_) in enumerate(zip(row, cw)):
        rect(s5, x, table_top + (ri + 1) * row_h, cw_, row_h, row_color)
        fc = GREEN if is_best else NAVY
        text(s5, x + 0.1, table_top + (ri + 1) * row_h + 0.18,
             cw_ - 0.1, row_h - 0.1,
             cell, 17, bold=is_best, color=fc, align=PP_ALIGN.CENTER)
        x += cw_

# delta 설명
text(s5, 0.5, 5.1, 12.5, 0.5,
     'NER Inject vs Baseline:  BLEU(c) +5.10  /  chrF +4.45  /  BLEU(m) +5.74  /  NER Recall +0.197',
     17, bold=True, color=GREEN)

# 막대 그래프 (BLEU(c))
text(s5, 0.5, 5.7, 5.0, 0.4, 'BLEU(c) 비교', 16, bold=True, color=NAVY)
bleu_vals = [35.51, 39.46, 39.78, 40.61]
bar_colors = [RGBColor(0x88,0x88,0x88), BLUE, RGBColor(0x15,0x6b,0x8a), GREEN]
bar_labels = ['Base', 'Few', '+fix', 'NER']
for i, (v, c, lb) in enumerate(zip(bleu_vals, bar_colors, bar_labels)):
    bar_w = (v / 50) * 5.5
    rect(s5, 0.5 + i * 3.1, 6.2, bar_w, 0.35, c)
    text(s5, 0.5 + i * 3.1, 6.1, 1.0, 0.25, lb, 13, color=NAVY)
    text(s5, 0.5 + i * 3.1 + bar_w + 0.05, 6.2, 1.0, 0.3, f'{v:.1f}', 13, color=NAVY)


# ─────────────────────────────────────────────
# 슬라이드 6: 결론
# ─────────────────────────────────────────────
s6 = add_slide()
bg(s6, NAVY)
rect(s6, 0, 0, 0.25, 7.5, GOLD)
rect(s6, 0, 5.8, 13.33, 1.7, BLUE)

text(s6, 0.5, 0.4, 12, 0.8, '결론 및 기여', 36, bold=True, color=GOLD)

conclusions = [
    ('NER Inject 효과 입증',
     'BLEU(c) 35.51 → 40.61 (+5.10), NER Recall 0.803 → 1.000'),
    ('인물명 정확도 완벽 달성',
     '356개 NER 엔티티 전량 올바르게 번역 (1.000)'),
    ('프롬프트 엔지니어링만으로 성능 향상',
     '파인튜닝 없이 few-shot + NER 주입으로 품질 개선'),
    ('승정원일기 전문 번역 가능성 확인',
     '실록 문체 보존 + 고유명사 정확도 동시 달성'),
]

for i, (title, desc) in enumerate(conclusions):
    rect(s6, 0.5, 1.3 + i * 1.05, 0.08, 0.5, GOLD)
    text(s6, 0.75, 1.3 + i * 1.05, 11.5, 0.4, title, 20, bold=True, color=WHITE)
    text(s6, 0.75, 1.72 + i * 1.05, 11.5, 0.35, desc, 16, color=RGBColor(0xaa, 0xcc, 0xff))

text(s6, 0.5, 6.05, 12, 0.5, '말모이 팀  |  승정원일기 자동 번역 프로젝트', 16, color=WHITE)


prs.save('말모이_승정원일기번역_발표.pptx')
print('저장 완료: 말모이_승정원일기번역_발표.pptx')
